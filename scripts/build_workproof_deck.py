"""Build the WorkProof deck from the Midnight Summit Template.

Strategy:
- Keep a curated subset of slides that preserves the template's
  multiple themes (light, grey, dark) and layout variety.
- Drop duplicate light/grey/dark twins so the deck does not repeat
  itself.
- Replace placeholder copy with WorkProof content while preserving
  the original run-level formatting (font, colour, size).
"""

from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt

SRC = "Midnight Summit Template Deck.pptx"
DST = "WorkProof Deck.pptx"

# (original_slide_index, content_dict)
# content_dict keys are shape indices on that slide.
# Each value is a list of strings; one string per paragraph.
# Empty strings clear a paragraph (used to drop placeholder bullets).
SLIDE_PLAN = [
    # 0 — Cover
    (0, {
        1: ["WorkProof"],
    }),
    # 6 — Light section header: the problem
    (6, {
        0: ["Every résumé is built on trust"],
    }),
    # 13 — Mission quote
    (13, {
        1: ["THE MISSION"],
        0: [
            "Replace trust with proof —",
            "without leaking a single byte of private data.",
        ],
    }),
    # 7 — Grey section header: how it works
    (7, {
        0: ["How WorkProof works"],
    }),
    # 9 — Light numbered content: three roles
    (9, {
        0: ["Three roles, one source of truth"],
        1: [
            "The cast",
            "Candidate submits a claim",
            "Identity hashed in-browser",
            "Verifier confirms a signal",
            "Recruiter sees a trust score",
            "Chain stores the commitment",
        ],
        2: ["01", "02", "03", "04", "05", "06"],
    }),
    # 8 — Dark section header: the AI
    (8, {
        0: ["The AI trust score"],
    }),
    # 15 — Three-column light: five signal channels (using 3 columns, two paired)
    (15, {
        1: ["Five signal channels"],
        0: [
            "Email & domain",
            "Hash-matched against the verifier's domain — never the raw address.",
            "",
            "",
            "",
            "",
        ],
        2: [
            "LinkedIn & duration",
            "Consistency of role and realistic tenure are scored together.",
            "",
            "",
            "",
            "",
        ],
        3: [
            "Endorsements & patterns",
            "Peer signals lift the score; suspicious patterns lower it.",
            "",
            "",
            "",
            "",
        ],
        4: ["01."],
        5: ["02."],
        6: ["03."],
    }),
    # 19 — Light two-column bullets: deterministic + anti-gameable
    (19, {
        5: ["Deterministic. Privacy-preserving. Anti-gameable."],
        3: ["What the AI is"],
        1: ["What the AI is not"],
        2: [
            "Deterministic — same inputs, same score",
            "Operates only on hashed data",
            "Explainable signal-by-signal breakdown",
            "Updates live as verifications arrive",
            "Open scoring rules, no black box",
            "Auditable trail per claim",
        ],
        0: [
            "Not a credit score or profile",
            "Not trained on personal data",
            "Cannot be inflated by spamming claims",
            "Cannot see emails or company names",
            "Cannot be bribed by a single verifier",
            "Cannot leak what it never received",
        ],
    }),
    # 21 — Image-left with bullets + intro: privacy by design
    (21, {
        3: ["Privacy by design"],
        4: [
            "Zero-knowledge proofs",
            "Claims are hashed in the browser before they leave the page. A ZK proof commits the claim on Midnight without revealing the candidate's identity or the exact dates worked.",
            "Recruiters see a trust score and the signals behind it — never an email, a real company name, or a résumé.",
            "",
            "",
            "",
            "",
        ],
        1: ["What stays private"],
        0: [
            "Email address",
            "Real company name",
            "Exact employment dates",
            "Verifier identity",
            "Raw résumé content",
            "Candidate's wallet address",
        ],
    }),
    # 11 — Dark numbered content: demo flow
    (11, {
        2: ["Live demo flow"],
        0: [
            "What you'll see",
            "Candidate submits a hashed claim",
            "ZK proof commits it on Midnight",
            "HR verifier confirms email domain",
            "Audit log updates in real time",
            "Trust score jumps +30 points",
            "Recruiter inspects — sees no PII",
        ],
        1: ["01", "02", "03", "04", "05", "06"],
    }),
    # 32 — Light single-column intro text: why Midnight
    (32, {
        1: ["Why Midnight"],
        0: [
            "The right substrate for verifiable trust",
            "Midnight gives us selective disclosure as a primitive — the chain proves what happened, the AI scores how much to trust it. ZK proofs are first-class, so private inputs never need to leave the user's device.",
            "The result: a public, tamper-proof record of verifications that reveals nothing about the people behind them. Trust and privacy stop being a trade-off.",
            "",
            "",
            "",
            "",
        ],
    }),
    # 36 — Light thank-you
    (36, {
        0: ["Thank you"],
        1: ["WorkProof Team\nBuilt on Midnight"],
        3: ["hello@workproof.xyz"],
    }),
]


def set_paragraph_text(paragraph, text):
    """Replace paragraph text, preserving the first run's formatting."""
    runs = paragraph.runs
    if runs:
        # Keep first run, drop the rest, then set its text.
        first = runs[0]
        first.text = text
        # Remove subsequent runs from the XML.
        r_elements = paragraph._p.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}r"
        )
        for r in r_elements[1:]:
            paragraph._p.remove(r)
    else:
        # No runs — write text directly via add_run on a new run-less paragraph.
        run = paragraph.add_run()
        run.text = text


def update_shape_text(shape, lines):
    """Replace text frame contents paragraph-by-paragraph."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    paragraphs = tf.paragraphs
    n_existing = len(paragraphs)
    n_new = len(lines)

    # Update existing paragraphs in place.
    for i in range(min(n_existing, n_new)):
        set_paragraph_text(paragraphs[i], lines[i])

    # If we have extra existing paragraphs, blank them out.
    for i in range(n_new, n_existing):
        set_paragraph_text(paragraphs[i], "")

    # If we need more paragraphs than existed, clone the last one's format.
    if n_new > n_existing and n_existing > 0:
        last_p = paragraphs[-1]._p
        for extra in lines[n_existing:]:
            new_p = deepcopy(last_p)
            last_p.addnext(new_p)
            last_p = new_p
        # Now refresh and set text on the new paragraphs.
        paragraphs = tf.paragraphs
        for i, line in enumerate(lines[n_existing:], start=n_existing):
            set_paragraph_text(paragraphs[i], line)


def main():
    prs = Presentation(SRC)
    all_slides = list(prs.slides)

    # Apply text updates first while indices still match.
    keep_indices = {idx for idx, _ in SLIDE_PLAN}
    plan_by_idx = {idx: content for idx, content in SLIDE_PLAN}

    for idx, content in plan_by_idx.items():
        slide = all_slides[idx]
        for shape_idx, lines in content.items():
            shapes = list(slide.shapes)
            if shape_idx >= len(shapes):
                print(f"  ! slide {idx}: shape {shape_idx} missing")
                continue
            update_shape_text(shapes[shape_idx], lines)

    # Now drop slides we don't want, in reverse order so indices stay stable.
    sldIdLst = prs.slides._sldIdLst
    sld_id_elems = list(sldIdLst)

    # Map keep indices to ordering (so we can also reorder to match SLIDE_PLAN).
    plan_order = [idx for idx, _ in SLIDE_PLAN]

    # Remove non-kept slides from the id list.
    for i, sld in enumerate(sld_id_elems):
        if i not in keep_indices:
            sldIdLst.remove(sld)
            # Also drop the relationship so the file stays clean.
            rId = sld.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass

    # Reorder the remaining sld id elements to match plan order.
    remaining = list(sldIdLst)
    # Map original-index -> element by matching rId back through prs.part.related_parts
    # Simpler: rebuild by iterating original list, picking those in keep order.
    kept_in_original_order = [
        sld_id_elems[i] for i in range(len(sld_id_elems)) if i in keep_indices
    ]
    # original_to_planpos
    plan_pos = {idx: pos for pos, idx in enumerate(plan_order)}
    kept_indices_sorted = sorted(
        [i for i in range(len(sld_id_elems)) if i in keep_indices],
        key=lambda i: plan_pos[i],
    )
    ordered_elems = [sld_id_elems[i] for i in kept_indices_sorted]

    # Clear and re-append in desired order.
    for el in remaining:
        sldIdLst.remove(el)
    for el in ordered_elems:
        sldIdLst.append(el)

    prs.save(DST)
    print(f"Saved {DST} with {len(ordered_elems)} slides.")


if __name__ == "__main__":
    main()
